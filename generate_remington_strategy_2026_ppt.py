"""Generate a 5-slide Remington Marketing Strategy 2026 deck with python-pptx.

Usage:
    python generate_remington_strategy_2026_ppt.py
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_FILE = "Remington_Marketing_Strategy_2026.pptx"

# Remington CI colors
RED = RGBColor(225, 6, 0)  # #E10600
BLACK = RGBColor(17, 17, 17)  # #111111
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
MID_GRAY = RGBColor(102, 102, 102)


def set_slide_bg(slide, color):
    bg_fill = slide.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = color


def add_header_bar(slide, title):
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.7),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RED
    bar.line.fill.background()

    tf = bar.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def add_body_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(11.8), Inches(0.9))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = BLACK


def add_bullets(slide, bullets, left=0.9, top=2.0, width=11.8, height=4.6):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    for idx, item in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)


def add_footer(slide, text="FREE TO HAIR. FREE TO BE YOU."):
    line = slide.shapes.add_shape(1, Inches(0.8), Inches(6.8), Inches(11.8), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(6.85), Inches(11.8), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = MID_GRAY
    p.alignment = PP_ALIGN.RIGHT


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank

    # Slide 1: Cover
    s1 = prs.slides.add_slide(layout)
    set_slide_bg(s1, WHITE)
    add_header_bar(s1, "Remington Marketing Strategy 2026")
    add_body_title(s1, "FREE TO HAIR.\nFREE TO BE YOU.")
    add_bullets(
        s1,
        [
            "Positioning: Affordable Premium That Lasts",
            "Focus: Hair Protection + Engineering Credibility + Durability",
            "Goal: Win the value-premium whitespace in Thailand",
        ],
        top=2.7,
        height=3.2,
    )
    add_footer(s1)

    # Slide 2: Target Segments
    s2 = prs.slides.add_slide(layout)
    set_slide_bg(s2, LIGHT_GRAY)
    add_header_bar(s2, "Target Audience Segmentation")
    add_body_title(s2, "4 Core Growth Segments")
    add_bullets(
        s2,
        [
            "1) First Jobbers (22-30): premium value, easy proof content",
            "2) Lifestyle Men (28-40+): expert-led grooming and design",
            "3) Modern Families (30-45): safety, reliability, long-term use",
            "4) Trend-Recreate Community: repeatable results and UGC proof loops",
        ],
    )
    add_footer(s2)

    # Slide 3: Product + AIRvive
    s3 = prs.slides.add_slide(layout)
    set_slide_bg(s3, WHITE)
    add_header_bar(s3, "Portfolio & Innovation")
    add_body_title(s3, "Core Range + AIRvive™ Launch (Jun 2026)")
    add_bullets(
        s3,
        [
            "Core price ladder: THB 3,500–4,500 (Keratin / PROluxe / Hydraluxe / Silk)",
            "AIRvive™ entry price from THB 4,900",
            "Hero SKUs: Digital Hair Dryer, Slim Straightener, 2-in-1 Air Styler, Rotating Wand",
            "Promise: Professional digital performance, simplified for daily users",
        ],
    )
    add_footer(s3)

    # Slide 4: Content & O2O
    s4 = prs.slides.add_slide(layout)
    set_slide_bg(s4, LIGHT_GRAY)
    add_header_bar(s4, "Content Engine & O2O Flywheel")
    add_body_title(s4, "Proof-Driven Growth Model")
    add_bullets(
        s4,
        [
            "Content mix: Product 60% | Lifestyle 15% | Engagement 10% | UGC/KOL 10% | Announcement 5%",
            "Short video formula (7-15s): 1 Pain → 1 Mechanism → 1 Result",
            "O2O cycle: Pop-up Demo → Social Share → PR Recap → UGC Loop → Conversion",
            "Trust drivers: Global heritage, visible proof, 2-year replacement clarity",
        ],
    )
    add_footer(s4)

    # Slide 5: Execution & KPI
    s5 = prs.slides.add_slide(layout)
    set_slide_bg(s5, WHITE)
    add_header_bar(s5, "2026 Execution Priorities")
    add_body_title(s5, "KPIs & Next Actions")
    add_bullets(
        s5,
        [
            "Launch AIRvive™ with creator seeding and always-on review cadence",
            "Scale weekly test-and-learn for short-form proof assets",
            "Track KPI: consideration lift, VDO completion, UGC volume, conversion, repeat rate",
            "Objective: Build Remington as Thailand's #1 Investment Piece Hair Tool Brand",
        ],
    )
    add_footer(s5)

    prs.save(OUTPUT_FILE)
    print(f"Created: {OUTPUT_FILE}")


if __name__ == "__main__":
    build_presentation()
